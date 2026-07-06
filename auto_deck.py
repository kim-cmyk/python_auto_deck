"""
auto_deck.py
================
Generate a 16:9 PowerPoint (.pptx) presentation from a JSON content file.

Usage
-----
    from auto_deck import build_presentation
    build_presentation("example_input.json", "output.pptx")

The input JSON describes one HEADER (global settings) and any number of PAGEs.
See ``example_input.json`` for a complete, valid example and ``README.md``
for the full format specification.

The loader is intentionally *tolerant*: it accepts // and /* */ comments,
Python literals (True/False/None), trailing commas, missing commas between
items, an un-wrapped top level, and repeated "PAGE" keys. Missing values fall
back to sensible defaults, so a partial / slightly-malformed file still builds.

Version : 1.0.0
Date    : 2026-07-06
"""

from __future__ import annotations

__version__ = "1.0.0"
__date__ = "2026-07-06"

import io
import json
import os
import re
from typing import Any, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

try:
    from PIL import Image
    _HAS_PIL = True
except Exception:  # pragma: no cover
    _HAS_PIL = False

try:
    import requests
    _HAS_REQUESTS = True
except Exception:  # pragma: no cover
    _HAS_REQUESTS = False


# ---------------------------------------------------------------------------
# Slide geometry (16:9 widescreen)
# ---------------------------------------------------------------------------
SLIDE_W = 13.333          # inches
SLIDE_H = 7.5             # inches
MARGIN = 0.7              # left / right page margin
HEADER_TOP = 0.35         # logo / header band top
TITLE_TOP = 0.55          # page title top on content slides
CONTENT_TOP = 1.7         # where the content area begins on content slides
FOOTER_H = 0.4
COL_GAP = 0.35            # gap between columns


# ---------------------------------------------------------------------------
# Themes  (3 light + 3 dark, all "formal / corporate")
# Every colour is a 6-digit hex string.
# ---------------------------------------------------------------------------
def _t(**kw) -> dict:
    return kw


THEMES = {
    # ---- LIGHT -----------------------------------------------------------
    "Bright1": _t(  # Corporate Blue
        bg="FFFFFF", title="1F3864", subtitle="2E5496", body="333333",
        accent="2E75B6", line="2E75B6",
        section_bg="1F3864", section_title="FFFFFF", section_sub="C9D6EA",
        header_text="1F3864", table_header_bg="1F3864", table_header_text="FFFFFF",
        table_band="EEF3FA", table_line="C9D6EA", footer="7F7F7F",
    ),
    "Bright2": _t(  # Teal / Green
        bg="FFFFFF", title="0B4F4A", subtitle="12726B", body="333333",
        accent="1FA39A", line="1FA39A",
        section_bg="0B4F4A", section_title="FFFFFF", section_sub="BFE3DF",
        header_text="0B4F4A", table_header_bg="0B4F4A", table_header_text="FFFFFF",
        table_band="E9F5F3", table_line="BFE3DF", footer="7F7F7F",
    ),
    "Bright3": _t(  # Warm Burgundy / Slate
        bg="FBFAF8", title="6B2737", subtitle="9E4A55", body="3A3A3A",
        accent="B08968", line="B08968",
        section_bg="6B2737", section_title="FFFFFF", section_sub="E6C9CE",
        header_text="6B2737", table_header_bg="6B2737", table_header_text="FFFFFF",
        table_band="F3EDE7", table_line="E3D8CE", footer="8A8A8A",
    ),
    # ---- DARK ------------------------------------------------------------
    "Dark1": _t(  # Navy
        bg="0E1B2A", title="FFFFFF", subtitle="AFC3DA", body="D5DEE8",
        accent="4C9BE0", line="4C9BE0",
        section_bg="09131F", section_title="FFFFFF", section_sub="AFC3DA",
        header_text="AFC3DA", table_header_bg="1B3A5B", table_header_text="FFFFFF",
        table_band="152438", table_line="2A4A5B", footer="8FA3B8",
    ),
    "Dark2": _t(  # Charcoal + Gold
        bg="1A1A1A", title="FFFFFF", subtitle="E8C98A", body="DDDDDD",
        accent="C9A227", line="C9A227",
        section_bg="0F0F0F", section_title="FFFFFF", section_sub="E8C98A",
        header_text="E8C98A", table_header_bg="3A3220", table_header_text="F5E6C0",
        table_band="262626", table_line="4A4A4A", footer="9A9A9A",
    ),
    "Dark3": _t(  # Deep Green
        bg="0F1E17", title="FFFFFF", subtitle="A9D3BE", body="D7E4DB",
        accent="3FA972", line="3FA972",
        section_bg="09140F", section_title="FFFFFF", section_sub="A9D3BE",
        header_text="A9D3BE", table_header_bg="16362A", table_header_text="FFFFFF",
        table_band="15271E", table_line="2A4A3A", footer="8FA79A",
    ),
}

DEFAULT_THEME = "Bright1"
DEFAULT_FONT_HEADINGS = "Calibri Light"
DEFAULT_FONT_BODY = "Calibri"

VALID_LAYOUTS = {
    "Title", "Section", "Single_Content", "Two_Content", "Tree_Content",
    "Four_Content", "Picture_Content", "Table_Content",
}

# how many content items each column layout renders
_COLUMN_COUNT = {"Single_Content": 1, "Two_Content": 2, "Tree_Content": 3,
                 "Four_Content": 4}


# ===========================================================================
#  Tolerant JSON loading
# ===========================================================================
class Obj(list):
    """A JSON object that preserves order and duplicate keys.

    Behaves like a list of (key, value) pairs, with dict-ish helpers.
    """

    def get(self, key: str, default: Any = None) -> Any:
        for k, v in self:
            if k == key:
                return v
        return default

    def getall(self, key: str) -> List[Any]:
        return [v for k, v in self if k == key]

    def has(self, key: str) -> bool:
        return any(k == key for k, _ in self)


# --- tokenizer used by the pre-processor -----------------------------------
_TOKEN_RE = re.compile(
    r"""
      (?P<ws>\s+)
    | (?P<str>"(?:\\.|[^"\\])*")
    | (?P<punct>[{}\[\]:,])
    | (?P<atom>[^\s{}\[\]:,"]+)
    """,
    re.VERBOSE,
)


def _strip_comments(text: str) -> str:
    """Remove // line comments and /* */ block comments (string-aware)."""
    out = []
    i, n = 0, len(text)
    in_str = False
    esc = False
    while i < n:
        c = text[i]
        if in_str:
            out.append(c)
            if esc:
                esc = False
            elif c == "\\":
                esc = True
            elif c == '"':
                in_str = False
            i += 1
            continue
        if c == '"':
            in_str = True
            out.append(c)
            i += 1
            continue
        if c == "/" and i + 1 < n and text[i + 1] == "/":
            i += 2
            while i < n and text[i] not in "\r\n":
                i += 1
            continue
        if c == "/" and i + 1 < n and text[i + 1] == "*":
            i += 2
            while i + 1 < n and not (text[i] == "*" and text[i + 1] == "/"):
                i += 1
            i += 2
            continue
        out.append(c)
        i += 1
    return "".join(out)


def _tokenize(text: str) -> List[Tuple[str, str]]:
    toks: List[Tuple[str, str]] = []
    pos = 0
    for m in _TOKEN_RE.finditer(text):
        if m.start() != pos:  # unexpected char, keep it as atom to avoid data loss
            toks.append(("atom", text[pos:m.start()]))
        kind = m.lastgroup
        toks.append((kind, m.group()))
        pos = m.end()
    if pos != len(text):
        toks.append(("atom", text[pos:]))
    return toks


_LITERAL_MAP = {"True": "true", "False": "false", "None": "null"}


def _repair(text: str) -> str:
    """Turn a lenient/JSON5-ish document into strict JSON text."""
    # normalise exotic whitespace that JSON rejects (nbsp, zero-width,
    # BOM, line/paragraph separators, ideographic space)
    for cp in (0xFEFF, 0x200B, 0x2060):
        text = text.replace(chr(cp), "")
    for cp in (0x00A0, 0x2028, 0x2029, 0x3000):
        text = text.replace(chr(cp), " ")
    text = _strip_comments(text)
    toks = _tokenize(text)

    # normalise Python literals (only bare atoms, never inside strings)
    toks = [(k, _LITERAL_MAP.get(v, v)) if k == "atom" else (k, v) for k, v in toks]

    # significant (non-whitespace) token indices
    sig = [i for i, (k, _) in enumerate(toks) if k != "ws"]

    def opens(idx: int) -> bool:
        k, v = toks[idx]
        return k in ("str", "atom") or v in ("{", "[")

    def closes(idx: int) -> bool:
        k, v = toks[idx]
        return k in ("str", "atom") or v in ("}", "]")

    insert_before = set()   # positions (token idx) that need a comma before them
    drop = set()            # comma token positions to delete (trailing commas)

    for a, b in zip(sig, sig[1:]):
        ka, va = toks[a]
        kb, vb = toks[b]
        # missing comma between two adjacent values
        if closes(a) and opens(b):
            insert_before.add(b)
        # trailing comma before a closing bracket
        if va == "," and vb in ("}", "]"):
            drop.add(a)

    # A brace group that contains no ':' is not a valid object -> treat it as
    # an array (e.g. a "set-literal" TABLE written as { "a","b","c" }).
    convert = {}            # token idx -> replacement bracket char
    stack = []
    for i, (k, v) in enumerate(toks):
        if k != "punct":
            continue
        if v in "{[":
            stack.append({"open": i, "colon": False, "char": v})
        elif v == ":" and stack:
            stack[-1]["colon"] = True
        elif v in "}]" and stack:
            fr = stack.pop()
            if fr["char"] == "{" and v == "}" and not fr["colon"]:
                convert[fr["open"]] = "["
                convert[i] = "]"

    out = []
    for i, (k, v) in enumerate(toks):
        if i in drop:
            continue
        if i in insert_before:
            out.append(",")
        out.append(convert.get(i, v))
    result = "".join(out)

    # wrap an un-wrapped top level:  "HEADER":{...},"PAGE":{...}  ->  { ... }
    first = next((v for (k, v) in ((tk, tv) for tk, tv in toks) if k != "ws"), "")
    if first.startswith('"'):
        result = "{" + result + "}"
    return result


def load_input(path_or_text: str) -> Obj:
    """Load a deck-definition file (or raw text) into an :class:`Obj` tree."""
    if os.path.exists(path_or_text):
        with open(path_or_text, "r", encoding="utf-8") as fh:
            raw = fh.read()
    else:
        raw = path_or_text

    try:
        return json.loads(raw, object_pairs_hook=Obj)
    except json.JSONDecodeError:
        repaired = _repair(raw)
        try:
            return json.loads(repaired, object_pairs_hook=Obj)
        except json.JSONDecodeError as exc:
            raise ValueError(
                "Could not parse the input file even after repair. "
                f"JSON error: {exc}. Please check for structural mistakes "
                "(a TABLE must be a 2-D array, e.g. [[\"h1\",\"h2\"],[\"a\",\"b\"]])."
            ) from exc


# ===========================================================================
#  Normalisation:  Obj tree  ->  clean python dicts w/ defaults
# ===========================================================================
def _norm_header(doc: Obj) -> dict:
    hdr = doc.get("HEADER")
    if not isinstance(hdr, Obj):
        hdr = Obj()
    theme = hdr.get("THEME", DEFAULT_THEME)
    if theme not in THEMES:
        theme = DEFAULT_THEME
    return {
        "logo": hdr.get("LOGO_LINK") or None,
        "theme": theme,
        "font_headings": hdr.get("FONT_HEADINGS") or DEFAULT_FONT_HEADINGS,
        "font_body": hdr.get("FONT_BODY") or DEFAULT_FONT_BODY,
        "page_number": bool(hdr.get("PAGE_NUMBER", True)),
    }


def _content_items(page: Obj) -> List[Tuple[str, Any]]:
    """Return ordered [(TYPE, value), ...] regardless of how content was written.

    Supports both the object form (possibly with duplicate keys) and the
    array-of-objects form.
    """
    pc = page.get("PAGE_CONTENT")
    items: List[Tuple[str, Any]] = []
    if isinstance(pc, Obj):
        items = [(k, v) for k, v in pc]
    elif isinstance(pc, list):
        for el in pc:
            if isinstance(el, Obj) and len(el):
                items.append((el[0][0], el[0][1]))
    return items


def _norm_pages(doc: Obj) -> List[dict]:
    raw_pages = doc.getall("PAGE")
    if not raw_pages:
        maybe = doc.get("PAGES")
        if isinstance(maybe, list):
            raw_pages = maybe

    pages = []
    for p in raw_pages:
        if not isinstance(p, Obj):
            continue
        layout = p.get("PAGE_LAYOUT", "Single_Content")
        if layout not in VALID_LAYOUTS:
            layout = "Single_Content"
        pages.append({
            "layout": layout,
            "title": p.get("PAGE_TITLE", "") or "",
            "subtitle": p.get("PAGE_SUBTITLE", "") or "",
            "content": _content_items(p),
        })
    return pages


# ===========================================================================
#  Low-level pptx drawing helpers
# ===========================================================================
def _rgb(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr)


def _set_bg(slide, hexstr: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hexstr)


def _add_rect(slide, left, top, width, height, hexstr):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(hexstr)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    return box, tf


def _style_run(run, font_name, size_pt, hexstr, bold=False, italic=False):
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(hexstr)


# ===========================================================================
#  Image fetching
# ===========================================================================
def _fetch_image(src: str) -> Optional[Tuple[io.BytesIO, Tuple[int, int]]]:
    """Return (stream, (w_px, h_px)) for a URL or local path, else None."""
    data: Optional[bytes] = None
    try:
        if re.match(r"^https?://", src, re.I):
            if not _HAS_REQUESTS:
                return None
            resp = requests.get(src, timeout=15, headers={"User-Agent": "auto_deck"})
            resp.raise_for_status()
            data = resp.content
        else:
            if not os.path.exists(src):
                return None
            with open(src, "rb") as fh:
                data = fh.read()
    except Exception:
        return None

    if data is None:
        return None

    # SVG is not supported by python-pptx; skip gracefully.
    head = data[:256].lstrip()
    if head[:5].lower() == b"<?xml" or head[:4].lower() == b"<svg":
        return None

    size = (800, 600)
    if _HAS_PIL:
        try:
            img = Image.open(io.BytesIO(data))
            size = img.size
            # convert unusual modes / formats pptx dislikes into PNG
            if img.format not in ("PNG", "JPEG", "GIF", "BMP", "TIFF"):
                buf = io.BytesIO()
                img.convert("RGB").save(buf, format="PNG")
                return buf, size
        except Exception:
            return None
    return io.BytesIO(data), size


def _place_image(slide, src, left, top, width, height, theme, font):
    """Fit an image inside a box, preserving aspect ratio & centring it."""
    result = _fetch_image(src)
    if result is None:
        _placeholder(slide, left, top, width, height, theme, font,
                     f"[image unavailable]")
        return
    stream, (w_px, h_px) = result
    box_ratio = width / height
    img_ratio = w_px / h_px if h_px else box_ratio
    if img_ratio > box_ratio:      # image wider -> fit width
        draw_w = width
        draw_h = width / img_ratio
    else:                           # image taller -> fit height
        draw_h = height
        draw_w = height * img_ratio
    off_l = left + (width - draw_w) / 2
    off_t = top + (height - draw_h) / 2
    try:
        slide.shapes.add_picture(stream, Inches(off_l), Inches(off_t),
                                 Inches(draw_w), Inches(draw_h))
    except Exception:
        _placeholder(slide, left, top, width, height, theme, font,
                     "[image could not be embedded]")


def _placeholder(slide, left, top, width, height, theme, font, text):
    _add_rect(slide, left, top, width, height, theme["table_band"])
    box, tf = _add_textbox(slide, left, top, width, height, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    _style_run(r, font["body"], 12, theme["footer"], italic=True)


# ===========================================================================
#  Content-item renderers
# ===========================================================================
def _render_bullets(tf, values, theme, font, size=16):
    if isinstance(values, str):
        values = [values]
    first = True
    for item in values:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(6)
        p.line_spacing = 1.1
        r = p.add_run()
        r.text = "•  " + str(item)
        _style_run(r, font["body"], size, theme["body"])


def _render_text(tf, value, theme, font, size=16, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.15
    r = p.add_run()
    r.text = str(value)
    _style_run(r, font["body"], size, theme["body"])


def _auto_body_size(items) -> int:
    """Very light heuristic so long text still fits."""
    total = 0
    for _t_, v in items:
        if isinstance(v, list):
            total += sum(len(str(x)) for x in v)
        else:
            total += len(str(v))
    if total > 1200:
        return 12
    if total > 700:
        return 14
    if total > 350:
        return 16
    return 18


def _render_table(slide, rows, left, top, width, height, theme, font):
    rows = [r for r in rows if isinstance(r, list)]
    if not rows:
        _placeholder(slide, left, top, width, height, theme, font, "[empty table]")
        return
    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    gtable = slide.shapes.add_table(n_rows, n_cols,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height)).table
    # size text depending on volume
    cell_size = 14 if n_rows <= 8 and n_cols <= 5 else 11
    for ci in range(n_cols):
        gtable.columns[ci].width = Inches(width / n_cols)
    for ri, row in enumerate(rows):
        for ci in range(n_cols):
            cell = gtable.cell(ri, ci)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            val = row[ci] if ci < len(row) else ""
            if ri == 0:  # header row
                cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(theme["table_header_bg"])
                colour = theme["table_header_text"]; bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(
                    theme["table_band"] if ri % 2 == 0 else theme["bg"])
                colour = theme["body"]; bold = False
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            _style_run(r, font["body"], cell_size, colour, bold=bold)


def _render_item(slide, kind, value, box_rect, theme, font, body_size):
    """Render one content item inside (left, top, width, height)."""
    left, top, width, height = box_rect
    kind = (kind or "").upper()
    if kind == "PICTURE":
        _place_image(slide, str(value), left, top, width, height, theme, font)
    elif kind == "TABLE":
        _render_table(slide, value if isinstance(value, list) else [],
                      left, top, width, height, theme, font)
    elif kind == "BULLET":
        _box, tf = _add_textbox(slide, left, top, width, height)
        _render_bullets(tf, value, theme, font, size=body_size)
    else:  # TEXT (default)
        _box, tf = _add_textbox(slide, left, top, width, height)
        _render_text(tf, value, theme, font, size=body_size)


# ===========================================================================
#  Chrome: logo, header title, footer
# ===========================================================================
class _Deck:
    def __init__(self, header: dict):
        self.header = header
        self.theme = THEMES[header["theme"]]
        self.font = {"head": header["font_headings"], "body": header["font_body"]}
        self.deck_title = ""     # "Title : Subtitle" of first Title page
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self.blank = self.prs.slide_layouts[6]
        self._logo_cache = None
        self._logo_tried = False

    # -- logo (top-left, small) -------------------------------------------
    def _logo(self):
        if self._logo_tried:
            return self._logo_cache
        self._logo_tried = True
        src = self.header["logo"]
        if src:
            self._logo_cache = _fetch_image(src)
        return self._logo_cache

    def add_logo(self, slide, dark_bg=False):
        logo = self._logo()
        if not logo:
            return
        stream, (w_px, h_px) = logo
        stream.seek(0)
        h = 0.4
        w = h * (w_px / h_px) if h_px else h * 3
        w = min(w, 2.2)
        try:
            slide.shapes.add_picture(stream, Inches(MARGIN), Inches(HEADER_TOP),
                                     height=Inches(h))
        except Exception:
            pass

    # -- footer ------------------------------------------------------------
    def add_footer(self, slide, page_no, on_light=True, show_title=True):
        theme, font = self.theme, self.font
        # bottom-left: deck title : subtitle
        if show_title and self.deck_title:
            box, tf = _add_textbox(slide, MARGIN, SLIDE_H - FOOTER_H - 0.15,
                                   SLIDE_W - 2 * MARGIN - 1.0, FOOTER_H,
                                   MSO_ANCHOR.MIDDLE)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = self.deck_title
            _style_run(r, font["body"], 9, theme["footer"])
        # bottom-right: page number
        if self.header["page_number"] and page_no is not None:
            box, tf = _add_textbox(slide, SLIDE_W - MARGIN - 1.0,
                                   SLIDE_H - FOOTER_H - 0.15, 1.0, FOOTER_H,
                                   MSO_ANCHOR.MIDDLE)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            r = p.add_run(); r.text = str(page_no)
            _style_run(r, font["body"], 9, theme["footer"])

    # -- page title (content slides) --------------------------------------
    def add_page_title(self, slide, title):
        theme, font = self.theme, self.font
        box, tf = _add_textbox(slide, MARGIN, TITLE_TOP,
                               SLIDE_W - 2 * MARGIN, 1.0, MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title
        _style_run(r, font["head"], 30, theme["title"], bold=True)
        # accent underline
        _add_rect(slide, MARGIN, TITLE_TOP + 1.02, 2.4, 0.05, theme["accent"])


# ===========================================================================
#  Slide builders (one per layout)
# ===========================================================================
def _build_title(deck: _Deck, page: dict, page_no):
    theme, font = deck.theme, deck.font
    slide = deck.prs.slides.add_slide(deck.blank)
    _set_bg(slide, theme["bg"])
    # left accent band
    _add_rect(slide, 0, 0, 0.35, SLIDE_H, theme["accent"])
    deck.add_logo(slide)

    box, tf = _add_textbox(slide, MARGIN + 0.3, 2.6,
                           SLIDE_W - 2 * MARGIN - 0.3, 1.6, MSO_ANCHOR.BOTTOM)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = page["title"]
    _style_run(r, font["head"], 46, theme["title"], bold=True)

    box2, tf2 = _add_textbox(slide, MARGIN + 0.3, 4.35,
                             SLIDE_W - 2 * MARGIN - 0.3, 1.0, MSO_ANCHOR.TOP)
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run(); r2.text = page["subtitle"]
    _style_run(r2, font["body"], 22, theme["subtitle"])
    # no page number / footer title on the very first title slide


def _build_section(deck: _Deck, page: dict, page_no):
    theme, font = deck.theme, deck.font
    slide = deck.prs.slides.add_slide(deck.blank)
    _set_bg(slide, theme["section_bg"])
    deck.add_logo(slide)

    box, tf = _add_textbox(slide, MARGIN, 2.9, SLIDE_W - 2 * MARGIN, 1.3,
                           MSO_ANCHOR.BOTTOM)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = page["title"]
    # a step "down" from the main title -> smaller & slightly softer weight
    _style_run(r, font["head"], 38, theme["section_title"], bold=True)

    _add_rect(slide, MARGIN, 4.25, 2.0, 0.05, theme["accent"])

    box2, tf2 = _add_textbox(slide, MARGIN, 4.4, SLIDE_W - 2 * MARGIN, 1.0)
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run(); r2.text = page["subtitle"]
    _style_run(r2, font["body"], 20, theme["section_sub"])

    deck.add_footer(slide, page_no, show_title=True)


def _build_columns(deck: _Deck, page: dict, page_no, n_cols):
    theme, font = deck.theme, deck.font
    slide = deck.prs.slides.add_slide(deck.blank)
    _set_bg(slide, theme["bg"])
    deck.add_logo(slide)
    deck.add_page_title(slide, page["title"])

    items = page["content"][:n_cols]      # excess is ignored
    body_size = _auto_body_size(items)

    avail_w = SLIDE_W - 2 * MARGIN
    col_w = (avail_w - COL_GAP * (n_cols - 1)) / n_cols
    top = CONTENT_TOP
    height = SLIDE_H - CONTENT_TOP - FOOTER_H - 0.2

    for i, (kind, value) in enumerate(items):
        left = MARGIN + i * (col_w + COL_GAP)
        # subtle card background for multi-column slides
        if n_cols > 1:
            _add_rect(slide, left, top, col_w, height, theme["table_band"])
            pad = 0.2
            box_rect = (left + pad, top + pad, col_w - 2 * pad, height - 2 * pad)
        else:
            box_rect = (left, top, col_w, height)
        _render_item(slide, kind, value, box_rect, theme, font, body_size)

    deck.add_footer(slide, page_no, show_title=True)


def _build_picture(deck: _Deck, page: dict, page_no):
    theme, font = deck.theme, deck.font
    slide = deck.prs.slides.add_slide(deck.blank)
    _set_bg(slide, theme["bg"])
    deck.add_logo(slide)
    deck.add_page_title(slide, page["title"])

    # keep only PICTURE + TEXT, ignore anything else / any excess
    picture = None
    text = None
    for kind, value in page["content"]:
        k = (kind or "").upper()
        if k == "PICTURE" and picture is None:
            picture = value
        elif k == "TEXT" and text is None:
            text = value

    top = CONTENT_TOP
    height = SLIDE_H - CONTENT_TOP - FOOTER_H - 0.2
    avail_w = SLIDE_W - 2 * MARGIN

    if text:  # picture left, small text right
        pic_w = avail_w * 0.58
        txt_w = avail_w * 0.42 - COL_GAP
        if picture:
            _place_image(slide, str(picture), MARGIN, top, pic_w, height,
                         theme, font)
        box, tf = _add_textbox(slide, MARGIN + pic_w + COL_GAP, top,
                               txt_w, height, MSO_ANCHOR.MIDDLE)
        _render_text(tf, text, theme, font, size=14)
    else:      # picture centred
        if picture:
            _place_image(slide, str(picture), MARGIN, top, avail_w, height,
                         theme, font)

    deck.add_footer(slide, page_no, show_title=True)


def _build_table(deck: _Deck, page: dict, page_no):
    theme, font = deck.theme, deck.font
    slide = deck.prs.slides.add_slide(deck.blank)
    _set_bg(slide, theme["bg"])
    deck.add_logo(slide)
    deck.add_page_title(slide, page["title"])

    table = None
    text = None
    for kind, value in page["content"]:
        k = (kind or "").upper()
        if k == "TABLE" and table is None:
            table = value
        elif k == "TEXT" and text is None:
            text = value

    top = CONTENT_TOP
    height = SLIDE_H - CONTENT_TOP - FOOTER_H - 0.2
    avail_w = SLIDE_W - 2 * MARGIN

    if text:  # table left, small text right
        tbl_w = avail_w * 0.62
        txt_w = avail_w * 0.38 - COL_GAP
        if table is not None:
            _render_table(slide, table if isinstance(table, list) else [],
                          MARGIN, top, tbl_w, height, theme, font)
        box, tf = _add_textbox(slide, MARGIN + tbl_w + COL_GAP, top,
                               txt_w, height, MSO_ANCHOR.MIDDLE)
        _render_text(tf, text, theme, font, size=14)
    else:      # table centred (use ~85% width)
        tbl_w = avail_w * 0.85
        left = MARGIN + (avail_w - tbl_w) / 2
        if table is not None:
            _render_table(slide, table if isinstance(table, list) else [],
                          left, top, tbl_w, height, theme, font)

    deck.add_footer(slide, page_no, show_title=True)


_BUILDERS = {
    "Title": _build_title,
    "Section": _build_section,
    "Picture_Content": _build_picture,
    "Table_Content": _build_table,
}


# ===========================================================================
#  Public entry point
# ===========================================================================
def build_presentation(input_path: str, output_path: str = "output.pptx") -> str:
    """Build a .pptx from a JSON deck definition and return the output path."""
    doc = load_input(input_path)
    header = _norm_header(doc)
    pages = _norm_pages(doc)

    deck = _Deck(header)

    # deck title (for footer) = first Title page's "title : subtitle"
    for pg in pages:
        if pg["layout"] == "Title":
            t = pg["title"].strip()
            s = pg["subtitle"].strip()
            deck.deck_title = f"{t} : {s}" if s else t
            break

    page_counter = 0
    for pg in pages:
        layout = pg["layout"]
        # numbering: everything except the opening Title slide gets a number
        if layout == "Title":
            page_no = None
        else:
            page_counter += 1
            page_no = page_counter

        if layout in _BUILDERS:
            _BUILDERS[layout](deck, pg, page_no)
        else:  # the four column layouts
            _build_columns(deck, pg, page_no, _COLUMN_COUNT.get(layout, 1))

    if not output_path.lower().endswith(".pptx"):
        output_path += ".pptx"
    deck.prs.save(output_path)
    return output_path


if __name__ == "__main__":
    import sys
    print(f"auto_deck v{__version__} ({__date__})")
    inp = sys.argv[1] if len(sys.argv) > 1 else "example_input.json"
    out = sys.argv[2] if len(sys.argv) > 2 else "output.pptx"
    path = build_presentation(inp, out)
    print(f"Saved presentation -> {path}")
